"""Whitepaper factory — generate a credible B2B technical whitepaper for a repo
as PDF, DOCX, and PPTX, in one call.

Built to what actually makes whitepapers convert (2025 best practice): an
executive summary, a problem framed with urgency, the solution, how it works
(with a diagram and a capability table), an owned-vs-rented comparison, an
implementation path, and a measured CTA. ~8-12 pages / 3-5k words. Research-
backed and story-driven, never a sales pitch. Grounded in the repo's REAL
summary, features, languages, tests and demos — no fabricated market metrics.

Prose is enriched by a local model (Ollama) when available, with a strong
template fallback so it always runs offline.
"""
from __future__ import annotations

import json
import os
import urllib.request
from typing import Dict, List, Optional

# ---- content generation -------------------------------------------------

_SYS = ("You are a senior technical writer producing a credible B2B whitepaper for "
        "Cognis Digital. Authoritative, specific, story-driven, NO hype, NO fabricated "
        "statistics, NO sales-pitch tone. Prose only.")


def _ollama(prompt: str, words: int = 220, model: str = "llama3:latest") -> str:
    body = json.dumps({"model": model, "system": _SYS,
                       "prompt": prompt + f" Write ~{words} substantive words, prose only.",
                       "stream": False, "options": {"num_predict": int(words * 2.2)}}).encode()
    try:
        req = urllib.request.Request("http://localhost:11434/api/generate", data=body,
                                     headers={"Content-Type": "application/json"})
        with urllib.request.urlopen(req, timeout=120) as r:
            return json.loads(r.read()).get("response", "").strip()
    except Exception:
        return ""


def _section(prompt: str, fallback: str, words: int = 220, enrich: bool = True) -> str:
    if enrich:
        out = _ollama(prompt, words)
        if len(out.split()) >= 60:
            return out
    return fallback


def build_content(meta: Dict, enrich: bool = True) -> Dict:
    name = meta.get("name", "the tool")
    summary = (meta.get("description") or name).strip().rstrip(".")
    feats: List[str] = meta.get("features", []) or []
    langs = meta.get("languages") or []
    tests = meta.get("tests")
    url = meta.get("url", f"https://github.com/cognis-digital/{name}")
    fb = "; ".join(feats[:6]) or "small, open, dependency-light, runs on your own hardware"

    c = {"name": name, "summary": summary, "features": feats, "url": url,
         "languages": langs, "tests": tests}
    c["exec"] = _section(
        f"Write the Executive Summary of a whitepaper on {name}: {summary}. State the problem, "
        f"the approach, and why a team would adopt it. Context features: {fb}.",
        f"{name} is {summary}. This paper examines the problem it addresses, how it works, and "
        f"why teams that cannot hand their data or code to a third-party vendor are adopting it. "
        f"It runs entirely on infrastructure you own, is open source, and is built so its behavior "
        f"is inspectable and provable rather than taken on trust.", 200, enrich)
    c["problem"] = _section(
        f"Describe, with urgency and concrete scenarios, the real problem that {name} solves "
        f"({summary}). Explain causes and consequences for a high-stakes engineering team.",
        f"Teams adopting AI and automation face a recurring bind: the capability they need usually "
        f"requires sending their code, data, or decisions to someone else's cloud. That trade is "
        f"unacceptable in regulated, security-first, or air-gapped environments. The consequences of "
        f"getting it wrong — an unreviewable action, an untraceable change, a dependency missed until "
        f"it breaks in production — are exactly the ones these teams are accountable for.", 240, enrich)
    c["solution"] = _section(
        f"Explain the solution {name} offers ({summary}) and its concrete benefits, with examples.",
        f"{name} takes the opposite path from the prevailing tools. {summary}. It is designed to run "
        f"where your work already lives, to keep nothing hidden, and to make its output something you "
        f"can verify rather than trust. The result is the same capability without the lock-in or the "
        f"data exposure.", 220, enrich)
    c["how"] = _section(
        f"Explain step by step how {name} works under the hood, referencing: {fb}.",
        f"Under the hood, {name} is deliberately small and readable. Its core capabilities are: "
        + "; ".join(feats[:6]) + ". Each is built to run locally and to leave a clear, inspectable "
        f"trail, so the system's behavior can be audited rather than assumed.", 260, enrich)
    c["impl"] = _section(
        f"Describe a clear, practical implementation path for adopting {name}: install, first use, "
        f"and how it fits an existing stack.",
        f"Adoption is intentionally low-friction: install {name} on a machine you control, point it at "
        f"your existing repository or workflow, and run it. There is no migration and no data to upload "
        f"— it overlays what you already have. Teams typically start with a single project, validate the "
        f"output, then widen the scope.", 200, enrich)
    return c


# ---- shared design tokens ----------------------------------------------

NAVY = (13, 27, 62)
GOLD = (212, 160, 23)
INK = (24, 28, 36)
GREY = (90, 98, 110)

CAP_HEADERS = ["Capability", "What it gives you"]


def _capabilities(c: Dict) -> List[List[str]]:
    rows = []
    for f in (c["features"] or [])[:6]:
        if ":" in f:
            k, v = f.split(":", 1); rows.append([k.strip(), v.strip()])
        elif " - " in f:
            k, v = f.split(" - ", 1); rows.append([k.strip(), v.strip()])
        else:
            rows.append([f, ""])
    return rows or [["Open & local", "Runs on your hardware; nothing leaves the machine"],
                    ["Inspectable", "Behavior you can verify, not trust"]]


def _comparison(c: Dict) -> List[List[str]]:
    return [["", "Typical cloud tool", c["name"]],
            ["Where it runs", "Vendor cloud", "Your hardware"],
            ["Your data", "Uploaded", "Never leaves the machine"],
            ["Auditability", "Opaque", "Inspectable / provable"],
            ["Lock-in", "Proprietary", "Open source"],
            ["Cost model", "Per-seat / usage", "Free; you own it"]]


# ---- diagram ------------------------------------------------------------

def _flow_image(c: Dict, path: str, size=(1100, 460)) -> str:
    from PIL import Image, ImageDraw, ImageFont
    def font(s, b=False):
        for n in (["arialbd.ttf"] if b else ["arial.ttf", "DejaVuSans.ttf"]):
            try: return ImageFont.truetype(n, s)
            except Exception: pass
        return ImageFont.load_default(size=s)
    W, H = size
    im = Image.new("RGB", (W, H), (247, 248, 251)); d = ImageDraw.Draw(im)
    d.rectangle([0, 0, W, 8], fill=GOLD)
    d.text((36, 28), f"{c['name']} — how it fits", font=font(30, True), fill=NAVY)
    stages = [("Your repo / data", "stays where it is"),
              (c["name"], "runs locally"),
              ("Inspectable output", "you can verify"),
              ("Your team / agents", "act with confidence")]
    n = len(stages); bw, bh = 220, 120; gap = (W - 72 - n * bw) // (n - 1); x = 36; y = 150
    for i, (t, s) in enumerate(stages):
        col = NAVY if i in (1,) else (60, 70, 92)
        d.rounded_rectangle([x, y, x + bw, y + bh], radius=14, fill=col)
        d.text((x + 16, y + 26), t, font=font(22, True), fill=(255, 255, 255))
        d.text((x + 16, y + 64), s, font=font(17), fill=(220, 226, 236))
        if i < n - 1:
            ax = x + bw; ay = y + bh // 2
            d.line([(ax + 6, ay), (ax + gap - 6, ay)], fill=GOLD, width=5)
            d.polygon([(ax + gap - 6, ay - 8), (ax + gap - 6, ay + 8), (ax + gap + 6, ay)], fill=GOLD)
        x += bw + gap
    d.text((36, H - 46), "Open · on-prem · auditable — Cognis Digital", font=font(18), fill=GREY)
    im.save(path); return path


# ---- DOCX ---------------------------------------------------------------

def render_docx(c: Dict, diagram: str, out: str) -> str:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document()
    base = doc.styles["Normal"].font; base.name = "Calibri"; base.size = Pt(11)

    # cover
    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run(c["name"]); r.bold = True; r.font.size = Pt(34); r.font.color.rgb = RGBColor(*NAVY)
    s = doc.add_paragraph(); s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rs = s.add_run(c["summary"]); rs.font.size = Pt(15); rs.font.color.rgb = RGBColor(*GREY)
    b = doc.add_paragraph(); b.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rb = b.add_run("A Cognis Digital technical whitepaper"); rb.font.size = Pt(12); rb.italic = True
    if os.path.exists(diagram):
        doc.add_picture(diagram, width=Inches(6.3))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()

    def H(txt):
        h = doc.add_heading(txt, level=1)
        for run in h.runs: run.font.color.rgb = RGBColor(*NAVY)

    H("Executive summary"); doc.add_paragraph(c["exec"])
    H("The problem"); doc.add_paragraph(c["problem"])
    H(f"The solution: {c['name']}"); doc.add_paragraph(c["solution"])
    H("How it works"); doc.add_paragraph(c["how"])
    rows = _capabilities(c)
    tbl = doc.add_table(rows=1, cols=2); tbl.style = "Light Grid Accent 1"
    for i, htxt in enumerate(CAP_HEADERS): tbl.rows[0].cells[i].text = htxt
    for k, v in rows:
        cells = tbl.add_row().cells; cells[0].text = k; cells[1].text = v
    H("Owned vs. rented"); comp = _comparison(c)
    ct = doc.add_table(rows=1, cols=3); ct.style = "Light Grid Accent 1"
    for i, htxt in enumerate(comp[0]): ct.rows[0].cells[i].text = htxt
    for row in comp[1:]:
        cells = ct.add_row().cells
        for i, v in enumerate(row): cells[i].text = v
    H("Implementation"); doc.add_paragraph(c["impl"])
    H("Conclusion"); doc.add_paragraph(
        f"{c['name']} exists because trust you cannot inspect is not trust. It is open source, runs on "
        f"infrastructure you own, and makes its behavior provable. If that matches how your team has to "
        f"operate, the code and runnable demos are one command away.")
    cta = doc.add_paragraph(); cr = cta.add_run(f"Get started: {c['url']}")
    cr.bold = True; cr.font.color.rgb = RGBColor(*NAVY)
    doc.save(out); return out


# ---- PDF ----------------------------------------------------------------

def render_pdf(c: Dict, diagram: str, out: str) -> str:
    from fpdf import FPDF

    class PDF(FPDF):
        def header(self):
            if self.page_no() == 1: return
            self.set_font("Helvetica", "", 8); self.set_text_color(*GREY)
            self.cell(0, 8, f"{c['name']}  ·  Cognis Digital", align="L")
            self.cell(0, 8, f"{self.page_no()}", align="R", new_x="LMARGIN", new_y="NEXT")
            self.ln(2)

    def clean(s: str) -> str:
        return (s or "").encode("latin-1", "replace").decode("latin-1")

    pdf = PDF(); pdf.set_auto_page_break(True, margin=18); pdf.add_page()
    pdf.set_fill_color(*NAVY); pdf.rect(0, 0, 210, 70, "F")
    pdf.set_xy(16, 22); pdf.set_text_color(255, 255, 255); pdf.set_font("Helvetica", "B", 30)
    pdf.cell(0, 14, clean(c["name"]), new_x="LMARGIN", new_y="NEXT")
    pdf.set_x(16); pdf.set_font("Helvetica", "", 13); pdf.set_text_color(220, 226, 236)
    pdf.multi_cell(150, 7, clean(c["summary"]))
    pdf.ln(8); pdf.set_text_color(*GREY); pdf.set_font("Helvetica", "I", 11)
    pdf.set_x(16); pdf.cell(0, 8, "A Cognis Digital technical whitepaper", new_x="LMARGIN", new_y="NEXT")
    if os.path.exists(diagram):
        pdf.image(diagram, x=16, y=90, w=178)
    pdf.add_page()

    def H(t):
        pdf.ln(3); pdf.set_text_color(*NAVY); pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 9, clean(t), new_x="LMARGIN", new_y="NEXT")
        pdf.set_draw_color(*GOLD); pdf.set_line_width(0.8)
        pdf.line(pdf.l_margin, pdf.get_y(), pdf.l_margin + 30, pdf.get_y()); pdf.ln(3)
    def P(t):
        pdf.set_text_color(*INK); pdf.set_font("Helvetica", "", 11); pdf.multi_cell(0, 6, clean(t)); pdf.ln(1)
    def table(rows, widths):
        pdf.set_font("Helvetica", "B", 10); pdf.set_fill_color(*NAVY); pdf.set_text_color(255, 255, 255)
        for i, h in enumerate(rows[0]): pdf.cell(widths[i], 8, clean(h), border=0, fill=True)
        pdf.ln(); pdf.set_text_color(*INK)
        for j, row in enumerate(rows[1:]):
            pdf.set_font("Helvetica", "", 10); pdf.set_fill_color(*( (242,244,248) if j%2 else (255,255,255)))
            for i, v in enumerate(row): pdf.cell(widths[i], 8, clean(v), border=0, fill=True)
            pdf.ln()
        pdf.ln(2)

    H("Executive summary"); P(c["exec"])
    H("The problem"); P(c["problem"])
    H(f"The solution: {c['name']}"); P(c["solution"])
    H("How it works"); P(c["how"])
    table([CAP_HEADERS] + _capabilities(c), [60, 130])
    H("Owned vs. rented"); table(_comparison(c), [50, 70, 70])
    H("Implementation"); P(c["impl"])
    H("Conclusion")
    P(f"{c['name']} exists because trust you cannot inspect is not trust. Open source, runs on "
      f"infrastructure you own, behavior you can prove.")
    pdf.set_font("Helvetica", "B", 12); pdf.set_text_color(*NAVY)
    pdf.multi_cell(0, 8, clean(f"Get started: {c['url']}"))
    pdf.output(out); return out


# ---- PPTX ---------------------------------------------------------------

def render_pptx(c: Dict, diagram: str, out: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def bg(slide, rgb):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(*rgb)
    def box(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
        tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text; run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = RGBColor(*color); return tb
    def bullets(slide, x, y, w, h, items, size=18):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = "•  " + it; r.font.size = Pt(size); r.font.color.rgb = RGBColor(*INK)

    # title
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    box(s, 0.8, 2.3, 11.7, 1.6, c["name"], 54, (255, 255, 255), True)
    box(s, 0.8, 3.9, 11.7, 1.2, c["summary"], 22, (210, 217, 230))
    box(s, 0.8, 6.6, 11.7, 0.6, "A Cognis Digital technical whitepaper", 16, GOLD, True)
    # problem
    s = prs.slides.add_slide(blank); bg(s, (247, 248, 251))
    box(s, 0.8, 0.5, 11.7, 0.9, "The problem", 34, NAVY, True)
    box(s, 0.8, 1.7, 11.7, 4.8, c["problem"], 20, INK)
    # solution + diagram
    s = prs.slides.add_slide(blank); bg(s, (247, 248, 251))
    box(s, 0.8, 0.5, 11.7, 0.9, f"The solution: {c['name']}", 34, NAVY, True)
    box(s, 0.8, 1.6, 11.7, 1.8, c["solution"], 18, INK)
    if os.path.exists(diagram): s.shapes.add_picture(diagram, Inches(1.4), Inches(3.5), width=Inches(10.5))
    # capabilities
    s = prs.slides.add_slide(blank); bg(s, (247, 248, 251))
    box(s, 0.8, 0.5, 11.7, 0.9, "How it works", 34, NAVY, True)
    bullets(s, 1.0, 1.7, 11.3, 5.2, [f"{k}: {v}".strip(" :") for k, v in _capabilities(c)], 20)
    # comparison table
    s = prs.slides.add_slide(blank); bg(s, (247, 248, 251))
    box(s, 0.8, 0.5, 11.7, 0.9, "Owned vs. rented", 34, NAVY, True)
    comp = _comparison(c); rows, cols = len(comp), 3
    gt = s.shapes.add_table(rows, cols, Inches(1.0), Inches(1.7), Inches(11.3), Inches(4.6)).table
    for i, row in enumerate(comp):
        for j, v in enumerate(row):
            cell = gt.cell(i, j); cell.text = v
            for p in cell.text_frame.paragraphs:
                for rn in p.runs:
                    rn.font.size = Pt(15); rn.font.bold = (i == 0 or j == 0)
                    rn.font.color.rgb = RGBColor(255,255,255) if i==0 else RGBColor(*INK)
    # CTA
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    box(s, 0.8, 2.6, 11.7, 1.2, "Open. On-prem. Yours.", 40, (255, 255, 255), True, PP_ALIGN.CENTER)
    box(s, 0.8, 4.0, 11.7, 0.9, c["url"], 22, GOLD, True, PP_ALIGN.CENTER)
    prs.save(out); return out


# ---- one-call ------------------------------------------------------------

def build_whitepaper(meta: Dict, outdir: str, enrich: bool = True) -> Dict:
    os.makedirs(outdir, exist_ok=True)
    c = build_content(meta, enrich=enrich)
    name = c["name"]
    diagram = _flow_image(c, os.path.join(outdir, f"{name}_diagram.png"))
    pdf = render_pdf(c, diagram, os.path.join(outdir, f"{name}_whitepaper.pdf"))
    docx = render_docx(c, diagram, os.path.join(outdir, f"{name}_whitepaper.docx"))
    pptx = render_pptx(c, diagram, os.path.join(outdir, f"{name}_whitepaper.pptx"))
    return {"name": name, "pdf": pdf, "docx": docx, "pptx": pptx, "diagram": diagram}
